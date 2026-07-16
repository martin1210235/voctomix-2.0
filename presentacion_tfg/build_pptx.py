#!/usr/bin/env python3
"""TFG defense deck (hybrid design) on the official ETSIT template.

Blue+cyan palette. Official layouts for cover / section dividers / closing;
custom visual content on the official master. Section dividers before every
part (Introducción, Arquitectura, Implementación, Resultados, Conclusiones).
Results analysed one by one with the memoria's own figures. Main deck + a
"Contenido adicional" backup section. Notes = spoken script (GUION_*.md).
"""
import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.oxml.ns import qn
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.chart.data import CategoryChartData
from PIL import Image

ROOT = "/home/sonda/Documentos/voctomix"
TPL = os.path.join(ROOT, "presentacion_tfg", "SFS17721.pptx")
OUT = os.path.join(ROOT, "presentacion_tfg", "presentacion_voctomix_etsit.pptx")

BLUE = RGBColor(0x08, 0x44, 0x7C)
CYAN = RGBColor(0x28, 0x99, 0xF9)
NAVY = RGBColor(0x02, 0x25, 0x40)
GRAY = RGBColor(0x46, 0x47, 0x47)
LGRAY = RGBColor(0x8A, 0x8C, 0x8E)
LINE = RGBColor(0xE0, 0xE4, 0xEA)
SOFT = RGBColor(0xEE, 0xF3, 0xF9)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)

FG = ROOT + "/presentacion_tfg/figuras/resultados/"
IMG = {
    "trad": ROOT + "/drawio_sources/produccion_tradicional.png",
    "arch": ROOT + "/drawio_sources/diagrama_sistema.png",
    "gui":  ROOT + "/figuras/gui_completa.png",
    "cn":   ROOT + "/figuras/voctogui_cybernemo1.png",
    "dock": ROOT + "/drawio_sources/escenario_docker.png",
    "k8s":  ROOT + "/drawio_sources/escenario_kubernetes.png",
    "s1":   ROOT + "/drawio_sources/escenario_1pc.png",
    "s2":   ROOT + "/drawio_sources/escenario_2pc.png",
    "s3":   ROOT + "/drawio_sources/escenario_docker.png",
    "s4":   ROOT + "/drawio_sources/escenario_kubernetes.png",
    "r_cpu": FG + "cpu.png", "r_ram": FG + "ram.png", "r_ene": FG + "energia.png",
    "r_lat": FG + "latencia.png", "r_tra": FG + "transicion.png",
    "r_estd": FG + "estab_docker.png", "r_estk": FG + "estab_k8s.png",
    "r_resd": FG + "resil_docker.png", "r_resk": FG + "resil_k8s.png",
}

prs = Presentation(TPL)
LAY = {l.name: l for l in prs.slide_layouts}
FOOTER = "Voctomix 2.0  ·  Defensa TFG  ·  ETSIT-UPM"


# ------------------------------------------------------------- primitives
def clear_slides(p):
    lst = p.slides._sldIdLst
    for sldId in list(lst):
        p.part.drop_rel(sldId.get(qn("r:id")))
        lst.remove(sldId)


def ph_by_idx(s, idx):
    for ph in s.placeholders:
        if ph.placeholder_format.idx == idx:
            return ph
    return None


def _put(tf, lines, align=PP_ALIGN.LEFT, anchor=None):
    tf.word_wrap = True
    if anchor:
        tf.vertical_anchor = anchor
    first = True
    for ln in lines:
        segs = ln if isinstance(ln, list) else [ln]
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.alignment = align
        for (text, size, bold, color) in segs:
            r = p.add_run()
            r.text = text; r.font.size = Pt(size); r.font.bold = bold
            r.font.name = "Arial"; r.font.color.rgb = color
    return tf


def textbox(s, l, t, w, h, lines, align=PP_ALIGN.LEFT, anchor=None):
    tb = s.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tb.text_frame.margin_left = 0; tb.text_frame.margin_right = 0
    _put(tb.text_frame, lines, align, anchor)
    return tb


def rrect(s, l, t, w, h, fill, line=None, radius=0.10):
    sh = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(l), Inches(t), Inches(w), Inches(h))
    try:
        sh.adjustments[0] = radius
    except Exception:
        pass
    if fill is None:
        sh.fill.background()
    else:
        sh.fill.solid(); sh.fill.fore_color.rgb = fill
    if line is None:
        sh.line.fill.background()
    else:
        sh.line.color.rgb = line; sh.line.width = Pt(1.25)
    sh.shadow.inherit = False
    return sh


def oval(s, l, t, d, fill):
    sh = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(l), Inches(t), Inches(d), Inches(d))
    sh.fill.solid(); sh.fill.fore_color.rgb = fill
    sh.line.fill.background(); sh.shadow.inherit = False
    return sh


def hline(s, l, t, w, color, weight=2.5):
    sh = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(l), Inches(t), Inches(w), Pt(weight))
    sh.fill.solid(); sh.fill.fore_color.rgb = color
    sh.line.fill.background(); sh.shadow.inherit = False
    return sh


def chip(s, l, t, w, text, color):
    sh = rrect(s, l, t, w, 0.32, color, radius=0.5)
    _put(sh.text_frame, [[(text, 10, True, WHITE)]], PP_ALIGN.CENTER, MSO_ANCHOR.MIDDLE)
    return sh


def add_image_fit(s, key, box):
    L, T, W, H = box
    path = IMG[key] if key in IMG else key
    w, h = Image.open(path).size
    ar = w / h
    if W / H > ar:
        nh, nw = H, H * ar
    else:
        nw, nh = W, W / ar
    return s.shapes.add_picture(path, Inches(L + (W - nw) / 2), Inches(T + (H - nh) / 2),
                                Inches(nw), Inches(nh))


def notes(s, text):
    s.notes_slide.notes_text_frame.text = text


def content(title, subtitle=None):
    s = prs.slides.add_slide(LAY["1_hoja sin elementos"])
    f = ph_by_idx(s, 36)
    if f is not None:
        _put(f.text_frame, [[(FOOTER, 9, False, LGRAY)]])
    textbox(s, 0.6, 0.92, 9.0, 0.62, [[(title, 25, True, BLUE)]])
    hline(s, 0.62, 1.55, 1.1, CYAN, 3)
    if subtitle:
        textbox(s, 0.6, 1.60, 8.8, 0.4, [[(subtitle, 13, False, LGRAY)]])
    return s


def section(title, subtitle):
    s = prs.slides.add_slide(LAY["Portadilla"])
    ph_by_idx(s, 12).text_frame.text = title
    ph_by_idx(s, 11).text_frame.text = subtitle
    d = ph_by_idx(s, 10)
    if d is not None:
        d.text_frame.text = "Defensa TFG · 2026"
    return s


def minicards(s, blocks, y=4.85):
    """3 mini-cards (tag + text) across the bottom."""
    x = 0.6
    for tag, txt, col in blocks:
        rrect(s, x, y, 2.85, 1.55, WHITE, line=LINE)
        t = rrect(s, x, y, 2.85, 0.42, col, radius=0.12)
        _put(t.text_frame, [[(tag, 11, True, WHITE)]], PP_ALIGN.CENTER, MSO_ANCHOR.MIDDLE)
        textbox(s, x + 0.2, y + 0.5, 2.5, 1.0, [[(txt, 11, False, GRAY)]])
        x += 3.0


def result(title, sub, figs, por_que, como, resultado, note):
    s = content(title, sub)
    if len(figs) == 1:
        add_image_fit(s, figs[0], (1.8, 1.95, 6.4, 2.7))
    else:
        add_image_fit(s, figs[0], (0.6, 1.95, 4.35, 2.7))
        add_image_fit(s, figs[1], (5.05, 1.95, 4.35, 2.7))
    minicards(s, [("POR QUÉ SE MIDE", por_que, CYAN), ("CÓMO", como, BLUE),
                  ("RESULTADO", resultado, BLUE)])
    notes(s, note)
    return s


def feature(title, sub, que, como, badge=None):
    s = content(title, sub)
    if badge:
        b = oval(s, 8.7, 0.85, 0.7, BLUE)
        _put(b.text_frame, [[(badge, 20, True, WHITE)]], PP_ALIGN.CENTER, MSO_ANCHOR.MIDDLE)
    rrect(s, 0.6, 2.1, 4.3, 3.2, WHITE, line=LINE)
    tag = rrect(s, 0.6, 2.1, 1.3, 0.45, CYAN, radius=0.2)
    _put(tag.text_frame, [[("QUÉ", 12, True, WHITE)]], PP_ALIGN.CENTER, MSO_ANCHOR.MIDDLE)
    textbox(s, 0.85, 2.7, 3.85, 2.5, [[(x, 12.5, False, GRAY)] for x in que])
    rrect(s, 5.1, 2.1, 4.3, 3.2, WHITE, line=LINE)
    tag2 = rrect(s, 5.1, 2.1, 1.3, 0.45, BLUE, radius=0.2)
    _put(tag2.text_frame, [[("CÓMO", 12, True, WHITE)]], PP_ALIGN.CENTER, MSO_ANCHOR.MIDDLE)
    textbox(s, 5.35, 2.7, 3.85, 2.5, [[(x, 12.5, False, GRAY)] for x in como])
    return s


def scenario(n, name, sub, imgkey, key_line, note):
    s = content(f"Escenario {n}: {name}", sub)
    o = oval(s, 8.7, 0.85, 0.7, BLUE)
    _put(o.text_frame, [[(str(n), 22, True, WHITE)]], PP_ALIGN.CENTER, MSO_ANCHOR.MIDDLE)
    add_image_fit(s, imgkey, (0.6, 2.0, 8.8, 3.3))
    bar = rrect(s, 0.6, 5.6, 8.8, 0.9, SOFT, radius=0.12)
    _put(bar.text_frame, [[(key_line, 13, False, GRAY)]], PP_ALIGN.CENTER, MSO_ANCHOR.MIDDLE)
    notes(s, note)
    return s


clear_slides(prs)

# =======================================================================
#  MAIN DECK
# =======================================================================

# 1. Portada
s = prs.slides.add_slide(LAY["PORTADA"])
ph_by_idx(s, 0).text_frame.text = "Producción remota de vídeo en tiempo real con software libre"
_put(ph_by_idx(s, 1).text_frame,
     [[("Voctomix 2.0", 16, True, WHITE)],
      [("Martín Herranz Sánchez  ·  Tutor: David Jiménez Bermejo", 12, False, WHITE)]])
notes(s, "[0:00] Buenos días. Presento mi TFG: producción remota de vídeo en tiempo real con "
         "software libre, el sistema Voctomix 2.0. En 15 minutos: el problema, la arquitectura, "
         "lo que he implementado, una demo en vídeo y los resultados.")

# 2. Índice
s = content("Índice")
items = ["Contexto y motivación", "Estado del arte", "Objetivos",
         "Arquitectura del sistema", "Contribuciones y demo",
         "Resultados y validación", "Impacto y méritos", "Conclusiones"]
for i, it in enumerate(items):
    col = 0.9 if i < 4 else 5.3
    yy = 2.1 + (i % 4) * 1.15
    oval(s, col, yy, 0.6, BLUE if i % 2 == 0 else CYAN)
    textbox(s, col, yy, 0.6, 0.6, [[(f"{i+1:02d}", 15, True, WHITE)]], PP_ALIGN.CENTER, MSO_ANCHOR.MIDDLE)
    textbox(s, col + 0.8, yy, 3.9, 0.6, [[(it, 14, False, GRAY)]], anchor=MSO_ANCHOR.MIDDLE)
notes(s, "[0:40] Este es el recorrido.")

# 3. Sección Introducción
section("Introducción", "El problema, el contexto y los objetivos")

# 4. El problema (visual: antes vs Voctomix)
s = content("El problema", "Producir un directo es caro y cerrado")
rrect(s, 0.6, 2.05, 4.2, 3.0, WHITE, line=LINE)
th = rrect(s, 0.6, 2.05, 4.2, 0.6, GRAY, radius=0.1)
_put(th.text_frame, [[("Producción tradicional", 14, True, WHITE)]], PP_ALIGN.CENTER, MSO_ANCHOR.MIDDLE)
textbox(s, 0.85, 2.8, 3.75, 2.2,
        [[("Mezclador hardware de miles de €", 13, False, GRAY)],
         [("Personal técnico especializado", 13, False, GRAY)],
         [("Dependencia de un fabricante", 13, False, GRAY)],
         [("Software (OBS/vMix): un operador,", 13, False, GRAY)],
         [("monolítico, sin API", 13, False, GRAY)]])
rrect(s, 5.2, 2.05, 4.2, 3.0, WHITE, line=CYAN)
th2 = rrect(s, 5.2, 2.05, 4.2, 0.6, BLUE, radius=0.1)
_put(th2.text_frame, [[("Voctomix 2.0", 14, True, WHITE)]], PP_ALIGN.CENTER, MSO_ANCHOR.MIDDLE)
textbox(s, 5.45, 2.8, 3.75, 2.2,
        [[("Software libre, 0 € de licencia", 13, True, BLUE)],
         [("Modular y contenerizado", 13, False, GRAY)],
         [("Control remoto y programable", 13, False, GRAY)],
         [("Se despliega con un comando", 13, False, GRAY)],
         [("Producción profesional accesible", 13, False, GRAY)]])
# arrow between
ar = s.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(4.8), Inches(3.25), Inches(0.42), Inches(0.6))
ar.fill.solid(); ar.fill.fore_color.rgb = CYAN; ar.line.fill.background(); ar.shadow.inherit = False
bar = rrect(s, 0.6, 5.35, 8.8, 0.95, SOFT, radius=0.12)
_put(bar.text_frame, [[("Democratizar la producción audiovisual", 14, True, BLUE),
                       ("   ·   contribuye a los ODS 9, 4 y 10 (innovación, educación, igualdad)",
                        12, False, GRAY)]], PP_ALIGN.CENTER, MSO_ANCHOR.MIDDLE)
notes(s, "[1:00] El problema es directo: producir un directo con varias cámaras cuesta miles de euros "
         "y necesita personal técnico. Voctomix 2.0 hace lo mismo con software libre, modular y "
         "contenerizado, a coste de licencia cero. En una palabra: democratiza la producción "
         "audiovisual. Por eso conecta con los ODS de la ONU: innovación, educación e igualdad.")

# 5. Contexto
s = content("Contexto y motivación", "Por qué hace falta")
cards = [("Coste elevado", "Mezcladores hardware de varios miles de euros."),
         ("Personal especializado", "Dependencia de técnicos y de un único fabricante."),
         ("Poca flexibilidad", "El software actual es de un solo operador, sin API.")]
x = 0.6
for t, d in cards:
    rrect(s, x, 2.15, 2.85, 2.2, WHITE, line=LINE)
    hline(s, x + 0.22, 2.5, 0.7, CYAN, 3)
    textbox(s, x + 0.22, 2.65, 2.45, 0.6, [[(t, 15, True, BLUE)]])
    textbox(s, x + 0.22, 3.25, 2.45, 1.0, [[(d, 12, False, GRAY)]])
    x += 3.0
add_image_fit(s, "trad", (0.9, 4.7, 8.2, 1.75))
notes(s, "[1:45] La cadena tradicional: cámaras, mezclador hardware, codificador y audiencia. "
         "Cara, cerrada e inaccesible para eventos de escala media. Ahí está la motivación.")

# 6. Estado del arte
s = content("Estado del arte", "Qué existía antes")
cols = [("Hardware dedicado", "Blackmagic ATEM: fiable pero muy costoso."),
        ("Software monolítico", "OBS, vMix: sin API ni modo headless."),
        ("Voctomix v1.3", "Cliente-servidor; usado en un evento real.")]
x = 0.6
for t, d in cols:
    rrect(s, x, 2.2, 2.85, 1.95, WHITE, line=LINE)
    top = rrect(s, x, 2.2, 2.85, 0.52, BLUE, radius=0.10)
    _put(top.text_frame, [[(t, 13, True, WHITE)]], PP_ALIGN.CENTER, MSO_ANCHOR.MIDDLE)
    textbox(s, x + 0.22, 2.9, 2.45, 1.1, [[(d, 12.5, False, GRAY)]])
    x += 3.0
textbox(s, 0.6, 4.55, 8.8, 1.4,
        [[("Modelo REMI: producción remota sobre IP", 15, True, CYAN)],
         [("Partimos de Voctomix v1.3 y lo llevamos a un sistema modular, contenerizado y con "
           "telemetría.", 13, False, GRAY)]])
notes(s, "[2:15] Hardware fiable pero caro; software potente pero cerrado. El modelo REMI plantea "
         "producir en remoto sobre IP. Parto de Voctomix v1.3 y lo llevo mucho más lejos.")

# 7. Objetivos
s = content("Objetivos", "Objetivo general: una plataforma libre, modular y reproducible de producción remota")
oes = [("OE-1", "Funcionalidades de producción profesional"),
       ("OE-2", "Servicio de telemetría en tiempo real"),
       ("OE-3", "Contenerización con Docker"),
       ("OE-4", "Despliegue en Kubernetes"),
       ("OE-5", "Validación en múltiples escenarios"),
       ("OE-6", "Despliegue en producción real")]
for i, (code, txt) in enumerate(oes):
    x = 0.6 + (i % 3) * 3.0
    yy = 2.15 + (i // 3) * 2.05
    rrect(s, x, yy, 2.85, 1.8, WHITE, line=LINE)
    badge = rrect(s, x + 0.22, yy + 0.25, 1.0, 0.5, BLUE if i % 2 == 0 else CYAN, radius=0.5)
    _put(badge.text_frame, [[(code, 13, True, WHITE)]], PP_ALIGN.CENTER, MSO_ANCHOR.MIDDLE)
    textbox(s, x + 0.22, yy + 0.9, 2.45, 0.8, [[(txt, 12.5, False, GRAY)]])
notes(s, "[2:30] Seis objetivos: funcionalidades profesionales, telemetría, Docker, Kubernetes, "
         "validación y despliegue real. Al final veréis que se cumplieron los seis.")

# 8. Sección Arquitectura
section("Arquitectura", "Cómo está construido el sistema")

# 9. Arquitectura
s = content("Arquitectura del sistema", "Cliente-servidor modular sobre TCP")
add_image_fit(s, "arch", (0.5, 1.95, 5.4, 4.6))
calls = [("voctocore", "Motor GStreamer: compone, mezcla y saca la señal"),
         ("voctogui", "Interfaz de control GTK 3"),
         ("TCP", "Separación estricta cliente-servidor"),
         ("RabbitMQ", "Telemetría desacoplada (AMQP)")]
y = 2.05
for i, (t, d) in enumerate(calls):
    oval(s, 6.05, y + 0.05, 0.28, BLUE if i % 2 == 0 else CYAN)
    textbox(s, 6.5, y, 3.1, 0.75, [[(t, 13, True, BLUE)], [(d, 11, False, GRAY)]])
    y += 0.95
xx = 6.05
for tech in ["GStreamer", "Python", "Docker"]:
    chip(s, xx, 5.95, 1.15, tech, BLUE); xx += 1.25
xx = 6.05
for tech in ["Kubernetes", "RabbitMQ"]:
    chip(s, xx, 6.35, 1.45, tech, CYAN); xx += 1.55
notes(s, "[3:05] Dos procesos: voctocore, el motor que compone vídeo, mezcla audio y saca la señal; "
         "y voctogui, la interfaz de control. Hablan por TCP. La telemetría va aparte por RabbitMQ. "
         "Todo software libre.")

# 10. Pipeline
s = content("Cadena de procesamiento", "De las fuentes a la señal de programa")
steps = [("Fuentes", "6 fuentes"), ("Compositor", "4 modos + transic."),
         ("Overlays", "logos y rótulos"), ("Blanker", "LIVE/PAUSE/NOSTR."),
         ("Salida", "programa y previews")]
hline(s, 0.9, 3.05, 8.2, LINE, 3)
x = 0.75
for i, (t, d) in enumerate(steps):
    o = oval(s, x + 0.55, 2.7, 0.7, BLUE if i % 2 == 0 else CYAN)
    _put(o.text_frame, [[(str(i + 1), 18, True, WHITE)]], PP_ALIGN.CENTER, MSO_ANCHOR.MIDDLE)
    textbox(s, x, 3.55, 1.75, 0.4, [[(t, 12.5, True, BLUE)]], PP_ALIGN.CENTER)
    textbox(s, x, 3.95, 1.75, 0.6, [[(d, 10.5, False, GRAY)]], PP_ALIGN.CENTER)
    x += 1.68
textbox(s, 0.6, 5.15, 8.8, 1.35,
        [[("Cada etapa es un módulo GStreamer independiente", 14, True, CYAN)],
         [("Añadí cada funcionalidad sin tocar el núcleo original de Voctomix.", 13, False, GRAY)],
         [("La salida de programa se codifica en H.264 y se distribuye por RTMP o HLS "
           "(YouTube, Twitch).", 12, False, LGRAY)]])
notes(s, "[3:30] El corazón es una tubería GStreamer. Las fuentes entran, el compositor aplica el "
         "modo, se ponen overlays, el blanker decide qué sale y se generan las salidas. Cada bloque "
         "es un módulo que añadí sin tocar el núcleo. La salida sale en H.264 por RTMP o HLS hacia "
         "plataformas como YouTube o Twitch.")

# 11. Sección Implementación
section("Implementación", "Lo que he desarrollado")

# 12. Contribuciones
s = content("Contribuciones", "Cinco bloques desarrollados")
conts = [("Composición y transiciones", "4 modos (FS, PiP, SbS, Lecture) con transiciones de 750 ms"),
         ("Overlays", "Logo PNG + 2 capas de texto, con fundido y auto-off en cada corte"),
         ("Stream Blanker", "Salida LIVE / PAUSE / NOSTREAM independiente de la mezcla"),
         ("Audio Follows Video", "Fundido cruzado automático del audio al cambiar de fuente"),
         ("Telemetría RabbitMQ", "Estado del mezclador en JSON vía AMQP, sin tocar el núcleo")]
for i, (t, d) in enumerate(conts):
    if i < 4:
        x = 0.6 + (i % 2) * 4.5; yy = 2.1 + (i // 2) * 1.55; w = 4.2
    else:
        x, yy, w = 0.6, 5.2, 8.7
    rrect(s, x, yy, w, 1.4, WHITE, line=LINE)
    rrect(s, x, yy, 0.18, 1.4, BLUE if i % 2 == 0 else CYAN, radius=0.02)
    num = oval(s, x + 0.35, yy + 0.42, 0.55, BLUE if i % 2 == 0 else CYAN)
    _put(num.text_frame, [[(str(i + 1), 16, True, WHITE)]], PP_ALIGN.CENTER, MSO_ANCHOR.MIDDLE)
    textbox(s, x + 1.1, yy + 0.2, w - 1.3, 0.45, [[(t, 14, True, BLUE)]])
    textbox(s, x + 1.1, yy + 0.68, w - 1.3, 0.65, [[(d, 11.5, False, GRAY)]])
notes(s, "[3:50] Cinco contribuciones: composición con cuatro modos y transiciones; overlays de logo "
         "y texto; el stream blanker de tres estados; el audio que sigue al vídeo; y la telemetría. "
         "Las vemos en directo en el vídeo. (El detalle de cómo funciona cada una está en el material "
         "de apoyo, diapositivas 34-38.)")

# 13. Video demo
s = content("Demostración en vídeo", "El recorrido por la herramienta (~4 min)")
add_image_fit(s, "gui", (0.55, 1.95, 5.7, 4.4))
oval(s, 2.95, 3.65, 1.0, WHITE)
tri = s.shapes.add_shape(MSO_SHAPE.ISOSCELES_TRIANGLE, Inches(3.25), Inches(3.9), Inches(0.5), Inches(0.5))
tri.rotation = 90
tri.fill.solid(); tri.fill.fore_color.rgb = BLUE; tri.line.fill.background(); tri.shadow.inherit = False
lbls = ["Composición y transiciones", "Overlays y rótulos", "Stream Blanker",
        "Audio Follows Video", "Telemetría en vivo"]
y = 2.05
for i, t in enumerate(lbls):
    rrect(s, 6.5, y, 3.0, 0.62, WHITE, line=CYAN if i % 2 else BLUE)
    hline(s, 6.5, y, 0.12, CYAN if i % 2 else BLUE, 18)
    textbox(s, 6.75, y, 2.7, 0.62, [[(t, 12, True, BLUE)]], anchor=MSO_ANCHOR.MIDDLE)
    y += 0.82
notes(s, "[4:10–8:10] NÚCLEO. Reproduzco el vídeo y narro cada función en directo. Ver GUION_VIDEO.md "
         "para el segundo a segundo. Llevo el mp4 como copia local además del embebido.")

# 14. Sección Resultados
section("Resultados y validación", "Cuatro escenarios y siete medidas")

# 15. Despliegue timeline
s = content("Despliegue", "Cuatro escenarios, de menos a más complejo")
tsteps = [("1", "Monopuesto", "Un equipo nativo"), ("2", "Dos PCs", "Red local"),
          ("3", "Docker Compose", "11 contenedores"), ("4", "Kubernetes", "Minikube, orquestado")]
hline(s, 1.0, 3.35, 8.0, LINE, 3)
x = 1.0
for i, (num, t, d) in enumerate(tsteps):
    o = oval(s, x + 0.9, 3.0, 0.7, BLUE if i % 2 == 0 else CYAN)
    _put(o.text_frame, [[(num, 20, True, WHITE)]], PP_ALIGN.CENTER, MSO_ANCHOR.MIDDLE)
    textbox(s, x, 3.9, 2.0, 0.5, [[(t, 14, True, BLUE)]], PP_ALIGN.CENTER)
    textbox(s, x, 4.35, 2.0, 0.6, [[(d, 11.5, False, GRAY)]], PP_ALIGN.CENTER)
    x += 2.0
textbox(s, 0.6, 5.4, 8.8, 1.0,
        [[("La misma funcionalidad escala de un portátil a un clúster ", 13, False, GRAY),
          ("sin tocar el código.", 13, True, BLUE)]])
notes(s, "[8:10] Lo validé en cuatro escenarios, de menos a más complejo. Los vemos uno a uno.")

# 16-19. Escenarios
scenario(1, "Monopuesto", "Todo en un solo equipo", "s1",
         "Todo en un portátil: fuentes, voctocore y voctogui. Sin red. El caso más simple.",
         "[8:25] Todo en un portátil. Las fuentes entran por los puertos 10000-10005, voctocore "
         "procesa, voctogui controla en local y el 15000 es lo que ve el espectador. Sin red.")
scenario(2, "Dos PCs", "Producción remota en red local", "s2",
         "voctocore en un PC servidor; el realizador controla desde otro PC por la LAN.",
         "[8:40] Separo motor y control: voctocore en un servidor, el realizador maneja voctogui "
         "desde otro PC por la red local. Ya es producción remota.")
scenario(3, "Docker Compose", "11 contenedores con un comando", "s3",
         "Las cámaras comparten la red de voctocore (localhost). Reproducible. Añade telemetría.",
         "[8:55] Todo en 11 contenedores que arrancan con un comando. Las cámaras comparten la red de "
         "voctocore. Aquí entra la telemetría a RabbitMQ. Reproducible en cualquier máquina.")
scenario(4, "Kubernetes", "Orquestación con Minikube", "s4",
         "El mismo sistema en un pod orquestado. RabbitMQ aparte. El realizador entra por NodePort.",
         "[9:10] El mismo sistema en Kubernetes: voctocore y cámaras en un pod, RabbitMQ aparte, el "
         "realizador entra por NodePort. Escala sin cambiar el código.")

# 20-22. Resultados clave en directo (CPU, latencia, resiliencia)
result("Resultados: CPU", "Consumo según el número de cámaras", ["r_cpu"],
       "Saber cuánto cuesta el sistema y si escala al añadir cámaras.",
       "Sesión de 5 min con 1-4 cámaras; mediana de la telemetría (cada 5 s).",
       "Docker: 33,5→38,8 % (+5,3 pp). K8s ~50 % fijo = Minikube. Amplio margen.",
       "[9:30] Consumo de CPU. Lo mido para ver si escala. Con cuatro cámaras solo llego al 38,8 %: "
       "muy eficiente. El 50 % fijo de Kubernetes es Minikube, no el mezclador.")
result("Resultados: latencia de conmutación", "Cuánto tarda un cambio de cámara", ["r_lat"],
       "En directo se cambia de cámara sin parar; un retraso se nota.",
       "30 cambios seguidos; tiempo desde la orden hasta que GStreamer lo aplica.",
       "Mediana 1,5 ms (Docker) / 1,8 ms (K8s). Umbral perceptible: 45 ms.",
       "[9:55] La latencia de cambio de cámara: 1,5 milisegundos de mediana. El umbral que se nota "
       "es 45. Estamos 30 veces por debajo: el cambio es instantáneo.")
result("Resultados: resiliencia", "Recuperación ante fallo de una cámara", ["r_resd", "r_resk"],
       "Si una cámara cae en directo, hay que recuperarla sola y rápido.",
       "Se mata el proceso de cam1 con kill -9; 10 iteraciones por entorno.",
       "Docker ~520 ms, K8s ~570 ms. Menos de 1 s, sin intervención.",
       "[10:20] Resiliencia: mato una cámara a propósito y mido cuánto tarda en volver. Medio "
       "segundo, sola, sin que yo haga nada. El espectador casi no lo nota. (RAM, energía, "
       "transición y estabilidad, en el material de apoyo.)")

# 27. Sección Conclusiones
section("Conclusiones", "Impacto, méritos y cierre")

# 28. Impacto y méritos (CyberNEMO + méritos fusionados)
s = content("Impacto y méritos", "Más allá de la memoria")
add_image_fit(s, "cn", (0.55, 2.0, 4.3, 3.0))
textbox(s, 0.55, 5.15, 4.3, 0.6,
        [[("Caso real: CyberNEMO", 13, True, BLUE),
          ("  ·  proyecto europeo de ciberseguridad (UPM)", 11, False, GRAY)]])
merits = [("Producción real", "Desplegado y usado en CyberNEMO, con participantes remotos."),
          ("Paper científico", "Artículo Voctomix 2.0 enviado a la revista MDPI Electronics."),
          ("Código abierto", "Sistema completo publicado como software libre en GitHub.")]
y = 2.0
for i, (t, d) in enumerate(merits):
    rrect(s, 5.1, y, 4.3, 1.35, WHITE, line=LINE)
    rrect(s, 5.1, y, 0.18, 1.35, BLUE if i % 2 == 0 else CYAN, radius=0.02)
    textbox(s, 5.45, y + 0.18, 3.8, 0.4, [[(t, 13.5, True, BLUE)]])
    textbox(s, 5.45, y + 0.62, 3.8, 0.6, [[(d, 11.5, False, GRAY)]])
    y += 1.5
notes(s, "[11:50] El trabajo no se quedó en la memoria: se ha usado en producción real en el proyecto "
         "europeo de ciberseguridad CyberNEMO, ha dado un artículo enviado a MDPI Electronics, y el "
         "código está publicado como software libre.")

# 29. Conclusiones
s = content("Conclusiones", "Los seis objetivos, cumplidos")
for i, (code, txt) in enumerate(oes):
    x = 0.6 + (i % 2) * 4.5
    yy = 2.15 + (i // 2) * 1.35
    rrect(s, x, yy, 4.2, 1.15, WHITE, line=LINE)
    ck = oval(s, x + 0.2, yy + 0.32, 0.5, BLUE if i % 2 == 0 else CYAN)
    _put(ck.text_frame, [[("✓", 18, True, WHITE)]], PP_ALIGN.CENTER, MSO_ANCHOR.MIDDLE)
    textbox(s, x + 0.9, yy + 0.18, 3.2, 0.8, [[(code + "  ", 12, True, BLUE), (txt, 11.5, False, GRAY)]],
            anchor=MSO_ANCHOR.MIDDLE)
textbox(s, 0.6, 6.35, 8.8, 0.6,
        [[("Producción broadcast profesional, con software libre y coste de licencia cero.", 13, True, BLUE)]],
        PP_ALIGN.CENTER)
notes(s, "[11:20] En conclusión, se cumplieron los seis objetivos. Voctomix 2.0 demuestra que se "
         "puede hacer producción broadcast profesional íntegramente con software libre y coste cero.")

# 29b. Líneas futuras
s = content("Líneas futuras", "Hacia dónde puede crecer Voctomix 2.0")
fut = [("Cámaras físicas", "Validar con tarjetas de captura en un directo real"),
       ("Salida H.265 / HEVC", "Mitad de ancho de banda a igual calidad"),
       ("Resolución 4K", "Escalar de 1080p a 4K en el pipeline de GStreamer"),
       ("GUI dinámica", "Añadir y quitar fuentes en caliente, sin reiniciar"),
       ("Más métricas", "Ampliar el análisis de rendimiento y fiabilidad")]
for i, (t, d) in enumerate(fut):
    yy = 2.05 + i * 0.92
    o = oval(s, 0.7, yy, 0.6, BLUE if i % 2 == 0 else CYAN)
    _put(o.text_frame, [[(str(i + 1), 16, True, WHITE)]], PP_ALIGN.CENTER, MSO_ANCHOR.MIDDLE)
    textbox(s, 1.5, yy + 0.02, 8.0, 0.55,
            [[(t + "   ", 15, True, BLUE), (d, 12.5, False, GRAY)]], anchor=MSO_ANCHOR.MIDDLE)
notes(s, "[11:40] Como líneas futuras: validar con cámaras físicas en un directo real; migrar la "
         "salida a H.265 para ahorrar ancho de banda; escalar a 4K; una interfaz que reconfigure las "
         "fuentes en caliente sin reiniciar; y ampliar las métricas.")

# 30. Cierre
s = prs.slides.add_slide(LAY["cierre de presentación"])
textbox(s, 0.6, 3.0, 6.0, 1.4,
        [[("Gracias por su atención", 30, True, WHITE)],
         [("Turno de preguntas", 17, False, RGBColor(0xD0, 0xE4, 0xFF))]])
notes(s, "[12:45] Gracias por su atención. Quedo a disposición del tribunal. Las siguientes "
         "diapositivas son material de apoyo.")

# =======================================================================
#  CONTENIDO ADICIONAL
# =======================================================================
section("Contenido adicional", "Material de apoyo para el turno de preguntas")

# Índice del contenido adicional
s = content("Contenido adicional", "Qué encontrarás aquí (para responder preguntas)")
menu = [("Mapa de puertos", "Interfaces de red del sistema"),
        ("Implementación", "Cómo funciona cada funcionalidad"),
        ("Ficheros nuevos", "Qué código se añadió y dónde"),
        ("Contenerización", "El truco de red de Docker"),
        ("Herramientas", "Cómo se midieron los resultados"),
        ("Kubernetes", "Los manifiestos de orquestación"),
        ("Resultados ampliados", "RAM, energía, transición, estabilidad"),
        ("Estado del arte", "Protocolos y códecs de contribución")]
for i, (t, d) in enumerate(menu):
    x = 0.6 + (i % 2) * 4.5
    yy = 2.05 + (i // 2) * 1.12
    rrect(s, x, yy, 4.2, 0.98, WHITE, line=LINE)
    tag = rrect(s, x + 0.15, yy + 0.28, 0.45, 0.45, BLUE if i % 2 == 0 else CYAN, radius=0.2)
    _put(tag.text_frame, [[(str(i + 1), 15, True, WHITE)]], PP_ALIGN.CENTER, MSO_ANCHOR.MIDDLE)
    textbox(s, x + 0.78, yy + 0.16, 3.3, 0.4, [[(t, 13.5, True, BLUE)]])
    textbox(s, x + 0.78, yy + 0.56, 3.3, 0.4, [[(d, 10.5, False, GRAY)]])
notes(s, "Material de apoyo, no se expone. Para responder al tribunal.")

# Puertos
s = content("Mapa de puertos del sistema", "Interfaces de red de voctocore y servicios")
ports = [("Puerto", "Servicio"),
         ("9999 / 9998", "Control (protocolo de comandos) y reloj de sincronización"),
         ("10000–10005", "Entrada de fuentes: cam1-4, break, intro"),
         ("11000 / 15000", "Salida de mezcla (raw) y salida de programa"),
         ("12000 / 14000–14005", "Previsualización de mezcla y de cada fuente (JPEG)"),
         ("17000 / 17001 / 18000", "Stream blanker (pause/offline) y audio de fondo (internos)"),
         ("8080", "Telemetría HTTP (estado del mezclador en JSON)"),
         ("5672 / 15672", "RabbitMQ: AMQP y panel de gestión web")]
tb = s.shapes.add_table(len(ports), 2, Inches(0.6), Inches(2.0), Inches(8.8), Inches(4.2)).table
tb.columns[0].width = Inches(2.6); tb.columns[1].width = Inches(6.2)
for r in range(len(ports)):
    for c in range(2):
        cell = tb.cell(r, c); cell.text = ports[r][c]
        p = cell.text_frame.paragraphs[0]; p.runs[0].font.size = Pt(12); p.runs[0].font.name = "Arial"
        if r == 0:
            p.runs[0].font.bold = True; p.runs[0].font.color.rgb = WHITE
            cell.fill.solid(); cell.fill.fore_color.rgb = BLUE
        else:
            p.runs[0].font.color.rgb = GRAY
            cell.fill.solid(); cell.fill.fore_color.rgb = WHITE if r % 2 else SOFT
notes(s, "El 15000 es la salida de programa: lo que ve el espectador. El 9999 el control. "
         "En Kubernetes se exponen como NodePorts.")

feature("Composición y transiciones", "Contribución 1 · composites.py, transitions.py, scene.py",
        ["4 modos: Full Screen, PiP,", "Side-by-Side y Lecture.",
         "Transiciones animadas de 750 ms", "entre cualquier par de modos."],
        ["Cada modo se define por geometría", "(posición, escala, alpha) en config.",
         "El motor interpola de la escena", "inicial a la final."], badge="1")
notes(prs.slides[-1], "Los 4 modos se definen en config; transitions.py interpola la escena en 750 ms.")

feature("Overlays: logos y rótulos", "Contribución 2 · overlay.py, toolbar/overlay.py",
        ["Logo PNG con transparencia.", "Dos capas de texto (lower-third).",
         "Rótulos independientes."],
        ["Fundido de 300 ms.", "Auto-off automático en cada corte.",
         "Se controla desde la GUI."], badge="2")
notes(prs.slides[-1], "overlay.py inyecta PNG y texto con fundido; se apaga solo en cada corte.")

feature("Stream Blanker", "Contribución 3 · streamblanker.py, toolbar/streamblank.py",
        ["Tres estados de salida:", "LIVE / PAUSE / NOSTREAM.",
         "La señal al público es distinta", "de la mezcla interna."],
        ["Fuentes TCP dedicadas para", "PAUSE y NOSTREAM.",
         "Conmutación limpia en frontera", "de buffer GStreamer."], badge="3")
notes(prs.slides[-1], "Separa lo que ve el público de lo que mezcla el realizador. Tres estados.")

s = content("Audio Follows Video y Telemetría", "Contribuciones 4 y 5")
rrect(s, 0.6, 2.05, 4.3, 3.6, WHITE, line=LINE)
tag = rrect(s, 0.6, 2.05, 2.3, 0.45, CYAN, radius=0.2)
_put(tag.text_frame, [[("Audio Follows Video", 11, True, WHITE)]], PP_ALIGN.CENTER, MSO_ANCHOR.MIDDLE)
textbox(s, 0.85, 2.65, 3.85, 2.9,
        [[("lib/audiomix.py", 12, True, BLUE)],
         [("Fundido cruzado automático del", 12, False, GRAY)],
         [("audio al cambiar de fuente.", 12, False, GRAY)], [(" ", 6, False, GRAY)],
         [("Quita la gestión manual del audio", 12, False, GRAY)],
         [("en directo. Menos carga para el", 12, False, GRAY)],
         [("realizador.", 12, False, GRAY)]])
rrect(s, 5.1, 2.05, 4.3, 3.6, WHITE, line=LINE)
tag2 = rrect(s, 5.1, 2.05, 1.8, 0.45, BLUE, radius=0.2)
_put(tag2.text_frame, [[("Telemetría", 11, True, WHITE)]], PP_ALIGN.CENTER, MSO_ANCHOR.MIDDLE)
textbox(s, 5.35, 2.65, 3.85, 2.9,
        [[("telemetry_service.py + gui_state_exporter.py", 10.5, True, BLUE)],
         [("Eventos CHANGE (al instante) y STATE", 12, False, GRAY)],
         [("(cada 5 s) en JSON Lines.", 12, False, GRAY)], [(" ", 6, False, GRAY)],
         [("Se publican en RabbitMQ por AMQP,", 12, False, GRAY)],
         [("sin tocar el núcleo. Permiten", 12, False, GRAY)],
         [("monitorizar y auditar desde fuera.", 12, False, GRAY)]])
notes(s, "AFV cruza el audio solo al cambiar de fuente. La telemetría publica eventos y estado en "
         "RabbitMQ por AMQP, sin tocar el núcleo.")

s = content("Ficheros incorporados", "Qué se añadió y dónde  ·  [N] = nuevo  ·  Anexo C")
tree = [
    "voctomix/",
    "├─ voctocore/lib/",
    "│    ├─ overlay.py  [N] logos y rótulos",
    "│    ├─ streamblanker.py  [N] LIVE / PAUSE / NOSTREAM",
    "│    ├─ transitions.py  [N] transiciones animadas",
    "│    └─ scene.py  [N] grafo de escena",
    "├─ voctogui/lib/",
    "│    ├─ gui_state_exporter.py  [N] estado → JSON",
    "│    └─ toolbar/",
    "│         ├─ overlay.py  [N] panel de overlays (GUI)",
    "│         └─ streamblank.py  [N] botones del blanker",
    "└─ example-scripts/ffmpeg/",
    "      └─ telemetry_service.py  [N] telemetría + RabbitMQ",
]
tb = textbox(s, 0.7, 2.0, 8.7, 3.8, [[(line, 11.5, False, GRAY)] for line in tree])
for p in tb.text_frame.paragraphs:
    for r in p.runs:
        r.font.name = "Consolas"
textbox(s, 0.7, 6.0, 8.7, 0.55,
        [[("+ infraestructura: docker-compose.yml · Dockerfile · k8s_escenario/ · tools/ · Makefile · 1pc/2pc",
           12, True, BLUE)]])
notes(s, "Resumen de ficheros incorporados. Si preguntan por código: overlay, blanker, transiciones, "
         "scene, GUI, telemetría y la infraestructura Docker/Kubernetes.")

s = content("Contenerización: la clave", "network_mode: service:voctocore")
add_image_fit(s, "dock", (0.6, 2.0, 5.4, 3.0))
textbox(s, 6.15, 2.05, 3.4, 4.4,
        [[("11 contenedores", 16, True, BLUE)],
         [("voctocore, 6 fuentes, stream", 12, False, GRAY)],
         [("blanker, audio, RabbitMQ y", 12, False, GRAY)],
         [("telemetría.", 12, False, GRAY)], [(" ", 8, False, GRAY)],
         [("Las cámaras comparten la red de", 12, False, GRAY)],
         [("voctocore y hablan por localhost,", 12, False, GRAY)],
         [("sin tocar una línea del código.", 12, False, GRAY)], [(" ", 8, False, GRAY)],
         [("La telemetría va en una red", 12, False, GRAY)],
         [("aparte: sin punto único de fallo.", 12, False, GRAY)]])
notes(s, "11 contenedores. La clave: network_mode service:voctocore, las cámaras comparten la red y "
         "hablan por localhost sin tocar el código.")

s = content("Herramientas de medición", "Cómo se obtuvieron los datos")
tools = [("measure_composite_latency.py", "Mide la latencia de cada cambio de fuente y de modo."),
         ("measure_resilience.py", "Mata el proceso de un contenedor y cronometra la recuperación."),
         ("analyze_stability.py", "Lee la telemetría (JSONL) y genera las curvas de CPU y RAM.")]
for i, (name, desc) in enumerate(tools):
    yy = 2.15 + i * 1.25
    rrect(s, 0.8, yy, 8.4, 0.95, WHITE, line=LINE)
    textbox(s, 1.05, yy + 0.18, 3.2, 0.35, [[(name, 12.5, True, BLUE)]])
    textbox(s, 4.45, yy + 0.18, 4.4, 0.4, [[(desc, 12, False, GRAY)]])
notes(s, "Herramientas reproducibles: medición de latencias, resiliencia y estabilidad desde JSONL.")

s = content("Despliegue en Kubernetes", "Orquestación declarativa")
add_image_fit(s, "k8s", (0.6, 2.0, 5.4, 3.0))
textbox(s, 6.15, 2.05, 3.4, 4.4,
        [[("Cinco manifiestos", 16, True, BLUE)],
         [("configmap, pvc, rabbitmq,", 12, False, GRAY)],
         [("studio y secret.", 12, False, GRAY)], [(" ", 8, False, GRAY)],
         [("Pod multi-contenedor con", 12, False, GRAY)],
         [("voctocore y las cámaras como", 12, False, GRAY)],
         [("sidecars.", 12, False, GRAY)], [(" ", 8, False, GRAY)],
         [("Sondas de preparación, arranque", 12, False, GRAY)],
         [("ordenado y volumen persistente.", 12, False, GRAY)]])
notes(s, "Cinco manifiestos declarativos: pod multi-contenedor, RabbitMQ aparte, sondas y volumen. "
         "Escala desde Docker sin tocar el código.")

# Resultados ampliados (movidos del directo)
result("Resultados: RAM", "Memoria según el número de cámaras", ["r_ram"],
       "Comprobar que la memoria no se dispara al añadir fuentes.",
       "Mismo test de 1-4 cámaras, medido por la telemetría.",
       "Plana: 9,7-10,1 % (Docker). voctocore reserva las 4 entradas al arrancar.",
       "La RAM es plana: voctocore reserva las entradas al arrancar; añadir cámaras casi no la mueve.")
result("Resultados: energía", "Consumo eléctrico neto (RAPL)", ["r_ene"],
       "Conocer el coste energético real atribuible a Voctomix.",
       "Sensor RAPL del procesador, descontando el consumo en reposo.",
       "25-32 W con 4 cámaras: menos del 20 % del TDP. Sobra margen.",
       "Consumo real medido con RAPL: 25-32 W, menos del 20 % de lo que aguanta el procesador.")
result("Resultados: latencia de transición", "Cambio de modo de composición", ["r_tra"],
       "Cambiar de modo (p. ej. a pantalla dividida) se anima frame a frame.",
       "32 medidas cambiando entre los 4 modos de composición.",
       "Docker 2,8 ms. K8s 79,8 ms (Minikube ocupa la CPU); en nube igualaría a Docker.",
       "Cambiar de modo se anima frame a frame: Docker 2,8 ms; K8s 80 ms por Minikube.")
result("Resultados: estabilidad", "31 minutos en marcha, sin degradación", ["r_estd", "r_estk"],
       "El CPU creep y las fugas de memoria solo salen a largo plazo.",
       "31 min con 4 cámaras; CPU y RAM cada 5 s (≈375 muestras).",
       "CPU estable, RAM totalmente plana, cero reinicios. Sin fugas.",
       "31 minutos seguidos: CPU estable y RAM plana. Ni fugas ni reinicios: apto para eventos largos.")

# Estado del arte: tecnologías de contribución (tablas de la memoria)
s = content("Estado del arte: tecnologías", "Protocolos y códecs de contribución audiovisual")


def mk_table(s, data, left, top, width, height, colws):
    tb = s.shapes.add_table(len(data), len(data[0]), Inches(left), Inches(top),
                            Inches(width), Inches(height)).table
    for ci, w in enumerate(colws):
        tb.columns[ci].width = Inches(w)
    for r in range(len(data)):
        for c in range(len(data[0])):
            cell = tb.cell(r, c); cell.text = data[r][c]
            p = cell.text_frame.paragraphs[0]; p.runs[0].font.size = Pt(10.5)
            p.runs[0].font.name = "Arial"
            if r == 0:
                p.runs[0].font.bold = True; p.runs[0].font.color.rgb = WHITE
                cell.fill.solid(); cell.fill.fore_color.rgb = BLUE
            else:
                p.runs[0].font.color.rgb = GRAY
                cell.fill.solid(); cell.fill.fore_color.rgb = WHITE if r % 2 else SOFT


textbox(s, 0.6, 1.9, 8.8, 0.35, [[("Protocolos", 12, True, CYAN)]])
mk_table(s, [("Protocolo", "Medio", "Latencia", "Uso típico"),
             ("SDI 12G", "Coaxial/fibra", "Sub-imagen", "Estudio broadcast"),
             ("SMPTE ST 2110", "UDP/IP", "Sub-imagen", "Producción IP nativa"),
             ("SRT", "UDP/IP (ARQ)", "< 1 s", "Contribución REMI"),
             ("RTP / RTSP", "UDP/IP", "Sub-imagen", "Cámaras IP")],
         0.6, 2.25, 8.8, 1.7, [2.2, 2.0, 1.6, 3.0])
textbox(s, 0.6, 4.1, 8.8, 0.35, [[("Códecs", 12, True, CYAN)]])
mk_table(s, [("Códec", "Compresión", "Latencia enc.", "Uso típico"),
             ("RAW I420 / PCM", "Sin compresión", "Nula", "Interno Voctomix"),
             ("Apple ProRes", "4:1 – 8:1", "< 1 frame", "Producción TV"),
             ("JPEG 2000", "4:1 – 20:1", "< 5 ms", "Contribución broadcast"),
             ("JPEG XS", "4:1 – 6:1", "< 1 ms", "Contribución 4K/8K")],
         0.6, 4.45, 8.8, 1.7, [2.2, 2.0, 1.6, 3.0])
textbox(s, 0.6, 6.3, 8.8, 0.5,
        [[("Voctomix usa RAW I420 + MKV/TCP: ", 11, True, BLUE),
          ("cero latencia de codificación en red local, frente al alto coste del SDI.",
           11, False, GRAY)]])
notes(s, "Contexto técnico para preguntas. En contribución hay dos familias: SDI (caro, cerrado) e "
         "IP (SMPTE 2110, SRT para REMI). En códecs, de RAW sin compresión a JPEG XS. Voctomix usa RAW "
         "I420 sobre MKV/TCP: cero latencia de codificación, ideal en red local.")

# Distribución de contenidos (salida al espectador)
s = content("Distribución de contenidos", "Cómo la señal de programa llega a la audiencia")
mk_table(s, [("Protocolo", "Latencia", "Uso principal"),
             ("RTMP", "2–5 s", "Ingesta a YouTube / Twitch / Facebook"),
             ("HLS", "2–30 s", "Reproducción masiva (CDN)"),
             ("MPEG-DASH", "2–30 s", "Streaming adaptativo abierto"),
             ("WebRTC", "0,1–0,5 s", "Monitorización y retornos")],
         0.6, 2.15, 8.8, 1.9, [2.4, 2.0, 4.4])
textbox(s, 0.6, 4.35, 8.8, 1.9,
        [[("La salida de Voctomix", 14, True, BLUE)],
         [("La señal de programa RAW I420 se codifica en H.264 con FFmpeg (libx264) en tiempo real,",
           12.5, False, GRAY)],
         [("se encapsula en MPEG-TS para RTMP o se segmenta en HLS. Compatible con cualquier CDN,",
           12.5, False, GRAY)],
         [("sin aceleración hardware. H.265/HEVC queda como línea futura para ahorrar ancho de banda.",
           12.5, False, GRAY)]])
notes(s, "La salida al público: voctocore saca RAW, FFmpeg lo codifica en H.264 en tiempo real y lo "
         "manda por RTMP a YouTube/Twitch o por HLS a una CDN. Todo con software libre.")

# Presupuesto económico
s = content("Presupuesto económico", "Anexo B  ·  coste de desarrollo del proyecto")
mk_table(s, [("Concepto", "Coste"),
             ("Mano de obra (320 h × 25 €)", "8.000,00 €"),
             ("Recursos materiales (amortización)", "1.000,00 €"),
             ("Gastos generales", "2.250,00 €"),
             ("IVA (21 %)", "2.504,25 €"),
             ("TOTAL", "14.429,25 €")],
         0.6, 2.15, 5.6, 2.7, [4.0, 1.6])
b = rrect(s, 6.5, 2.4, 3.0, 1.6, BLUE)
_put(b.text_frame, [[("0 €", 34, True, WHITE)], [("en licencias de software", 12, False, WHITE)]],
     PP_ALIGN.CENTER, MSO_ANCHOR.MIDDLE)
textbox(s, 6.5, 4.15, 3.1, 2.0,
        [[("Todo el sistema es software libre.", 12.5, True, BLUE)],
         [("El coste es desarrollo y hardware,", 12, False, GRAY)],
         [("no licencias, frente a los miles de", 12, False, GRAY)],
         [("euros de un mezclador comercial.", 12, False, GRAY)]])
notes(s, "El presupuesto del TFG: unos 14.400 €, casi todo mano de obra y equipos. Lo importante para "
         "el argumento: 0 € en licencias, porque todo es software libre, frente al hardware de miles de €.")

prs.core_properties.title = "Producción remota de vídeo en tiempo real con software libre"
prs.core_properties.author = "Martín Herranz Sánchez"
prs.save(OUT)
print("Saved:", OUT, "slides:", len(prs.slides._sldIdLst))
