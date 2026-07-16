#!/usr/bin/env python3
"""Comprobaciones automáticas de la presentación (formato y coherencia).

Uso:  python3 verificar.py
Revisa lo mecanizable del PRESENTACION_CHECKLIST.md.
"""
import sys
import zipfile
from pptx import Presentation
from pptx.util import Emu

PPTX = "presentacion_voctomix_etsit.pptx"
MAIN_LAST = 29          # las diapos 1..29 son el cuerpo principal
EXPECTED_MIN = 45       # nº de diapos esperado (aviso si baja de esto)

ok, warn, err = [], [], []


def main():
    # 1) duplicados de parte slideXX.xml
    with zipfile.ZipFile(PPTX) as z:
        names = [n for n in z.namelist() if n.startswith("ppt/slides/slide") and n.endswith(".xml")]
    dups = {n for n in names if names.count(n) > 1}
    (err if dups else ok).append(f"partes de diapo duplicadas: {sorted(dups) or 'ninguna'}")

    p = Presentation(PPTX)
    n = len(p.slides._sldIdLst)
    (ok if n >= EXPECTED_MIN else warn).append(f"nº de diapositivas: {n}")

    # 2) placeholders sueltos en el cuerpo principal
    leftovers = []
    for i, s in enumerate(p.slides, start=1):
        for sh in s.shapes:
            if sh.has_text_frame:
                t = sh.text_frame.text.lower()
                if ("insertar foto" in t or "insertar vídeo" in t or "insertar video" in t):
                    leftovers.append((i, sh.text_frame.text.strip()[:25]))
    main_left = [x for x in leftovers if x[0] <= MAIN_LAST]
    (err if main_left else ok).append(f"placeholders en cuerpo principal (1-{MAIN_LAST}): {main_left or 'ninguno'}")
    if any(x[0] > MAIN_LAST for x in leftovers):
        warn.append(f"placeholders en contenido adicional (tolerado): {[x for x in leftovers if x[0] > MAIN_LAST]}")

    # 3) bloque de título en diapos de contenido (título ~y=0.92 + pie)
    sin_titulo = []
    for i, s in enumerate(p.slides, start=1):
        tiene_titulo = any(sh.has_text_frame and sh.top is not None and 0.7 < Emu(sh.top).inches < 1.2
                           and sh.text_frame.text.strip() and not sh.text_frame.text.startswith("Defensa")
                           for sh in s.shapes)
        es_divisor = any(sh.has_text_frame and sh.text_frame.text.strip().isupper()
                         and len(sh.text_frame.text.strip()) > 4 for sh in s.shapes)
        if not tiene_titulo and not es_divisor and i not in (1, 2, 29):
            sin_titulo.append(i)
    (ok if not sin_titulo else warn).append(f"diapos sin bloque de título claro: {sin_titulo or 'ninguna'}")

    # informe
    print("=" * 60)
    for m in ok:
        print("  OK   ", m)
    for m in warn:
        print("  AVISO", m)
    for m in err:
        print("  ERROR", m)
    print("=" * 60)
    if err:
        print("RESULTADO: hay ERRORES que corregir.")
        sys.exit(1)
    print("RESULTADO: sin errores." + (" Revisa los avisos." if warn else ""))


if __name__ == "__main__":
    main()
